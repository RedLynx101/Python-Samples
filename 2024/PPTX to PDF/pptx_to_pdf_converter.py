import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image
import io
import os
from reportlab.lib.utils import ImageReader
from pptx.enum.shapes import MSO_SHAPE_TYPE

class PPTXtoPDFConverter:
    def __init__(self, master):
        self.master = master
        master.title("PPTX to PDF Converter")
        master.geometry("400x200")

        self.label = tk.Label(master, text="Select PPTX files to convert:")
        self.label.pack(pady=10)

        self.select_button = tk.Button(master, text="Select Files", command=self.select_files)
        self.select_button.pack(pady=10)

        self.convert_button = tk.Button(master, text="Convert to PDF", command=self.convert_to_pdf)
        self.convert_button.pack(pady=10)

        self.files = []

    def select_files(self):
        self.files = filedialog.askopenfilenames(filetypes=[("PowerPoint files", "*.pptx")])
        if self.files:
            self.label.config(text=f"{len(self.files)} file(s) selected")
        else:
            self.label.config(text="No files selected")

    def convert_to_pdf(self):
        if not self.files:
            messagebox.showerror("Error", "No files selected")
            return

        output_file = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if not output_file:
            return

        c = canvas.Canvas(output_file, pagesize=letter)

        for pptx_file in self.files:
            prs = Presentation(pptx_file)
            for slide in prs.slides:
                # Start a new page for each slide
                c.showPage()
                
                # Set initial text position
                text_y_position = 750  # Starting Y position for text

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text = run.text
                                c.drawString(50, text_y_position, text)
                                text_y_position -= 15  # Move down for the next line of text

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_stream = io.BytesIO(shape.image.blob)
                        img = Image.open(img_stream)
                        img_width, img_height = img.size
                        pdf_width, pdf_height = letter

                        scale = min(pdf_width / img_width, pdf_height / img_height)
                        new_width = img_width * scale
                        new_height = img_height * scale

                        x = (pdf_width - new_width) / 2
                        y = (pdf_height - new_height) / 2

                        c.drawImage(ImageReader(img), x, y, width=new_width, height=new_height)

        c.save()
        messagebox.showinfo("Success", f"PDF saved as {output_file}")

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTXtoPDFConverter(root)
    root.mainloop()
